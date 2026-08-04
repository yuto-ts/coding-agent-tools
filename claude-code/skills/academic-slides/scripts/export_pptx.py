#!/usr/bin/env python3
"""レイアウトJSONから編集可能な .pptx を生成する(Google Slides 取り込み用)。

  bash export_pptx.sh deck.html deck.pptx   # 通常はこちらを使う
  python3 export_pptx.py layout.json out.pptx

テキストは Google Slides 上で編集できる本物のテキストボックスとして出力する。
フォントは Google Slides が標準で持つ Noto Sans JP を指定する。英語のみの
デッキなら FONT_JP を Roboto や Arial に変えてよい。

既知の制約(ユーザーに事前に伝えること):
  - フォント置換で行の折り返しが数文字ずれる
  - 破線は実線になる
  - 1行が1テキストボックスになるため、段落をまたぐ一括編集はできない
"""
import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

BASE = os.getcwd()  # 相対パスの画像を解決する起点(build() で layout.json の位置に更新)
PX = 9525  # 1px = 9525 EMU (96dpi)
FONT_JP = "Noto Sans JP"
FONT_SERIF = "Noto Serif JP"

ALIGN = {
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "left": PP_ALIGN.LEFT,
    "start": PP_ALIGN.LEFT,
    "end": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def emu(px):
    return Emu(int(round(px * PX)))


def color(c):
    return RGBColor(c[0], c[1], c[2])


def set_font(run, name):
    """latin / 東アジア / complex script すべてに同じフォントを指定する。"""
    run.font.name = name
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        for old in rpr.findall(qn(tag)):
            rpr.remove(old)
        el = rpr.makeelement(qn(tag), {"typeface": name})
        rpr.append(el)


def no_shadow(shape):
    sppr = shape._element.spPr
    for tag in ("a:effectLst",):
        for old in sppr.findall(qn(tag)):
            sppr.remove(old)
    sppr.append(sppr.makeelement(qn("a:effectLst"), {}))


def add_rect(slide, s):
    radius = s.get("radius") or 0
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius >= 4 else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, emu(s["x"]), emu(s["y"]), emu(s["w"]), emu(s["h"]))
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            shp.adjustments[0] = min(0.5, radius / max(1.0, min(s["w"], s["h"])))
        except (IndexError, ValueError):
            pass
    if s.get("fill"):
        shp.fill.solid()
        shp.fill.fore_color.rgb = color(s["fill"])
    else:
        shp.fill.background()
    line = s.get("line")
    if line and line.get("c"):
        shp.line.color.rgb = color(line["c"])
        shp.line.width = Pt(line["w"] * 0.75)
    else:
        shp.line.fill.background()
    no_shadow(shp)
    shp.text_frame.text = ""
    return shp


def add_image(slide, s):
    # export_pptx.sh が src を絶対パス化しているので通常はそのまま解決できる。
    # 相対パスで呼ばれた場合は layout.json のあるディレクトリを起点にする。
    path = s["src"] if os.path.isabs(s["src"]) else os.path.join(BASE, s["src"])
    if not os.path.exists(path):
        print(f"  ! 画像が見つかりません: {s['src']}", file=sys.stderr)
        return None
    return slide.shapes.add_picture(path, emu(s["x"]), emu(s["y"]), emu(s["w"]), emu(s["h"]))


def add_text(slide, s):
    # 高さは余裕を持たせる(フォント差で溢れても隠れないように)
    tb = slide.shapes.add_textbox(emu(s["x"]), emu(s["y"] - 2), emu(s["w"]), emu(s["h"] + 6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    # <br> でパラグラフを分割
    paras = [[]]
    for r in s["runs"]:
        if r.get("br"):
            paras.append([])
        else:
            paras[-1].append(r)

    first = True
    for runs in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = ALIGN.get(s.get("align", "left"), PP_ALIGN.LEFT)
        p.line_spacing = round(s.get("lh", 1.4), 2)
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        if s.get("marL") is not None:  # 箇条書きのぶら下げインデント
            ppr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
            ppr.set("marL", str(int(round(s["marL"] * PX))))
            ppr.set("indent", str(int(round(s.get("indent", 0) * PX))))
        for r in runs:
            run = p.add_run()
            run.text = r["t"]
            run.font.size = Pt(round(r["size"] * 0.75, 1))
            run.font.bold = bool(r.get("bold"))
            run.font.italic = bool(r.get("italic"))
            run.font.color.rgb = color(r.get("color", [0, 0, 0]))
            set_font(run, FONT_SERIF if r.get("serif") else FONT_JP)
            if r.get("sub") or r.get("sup"):
                run.font._rPr.set("baseline", "30000" if r.get("sup") else "-25000")
    return tb


def build(layout_path, out_path):
    global BASE
    BASE = os.path.dirname(os.path.abspath(layout_path))
    with open(layout_path, encoding="utf-8") as f:
        pages = json.load(f)

    prs = Presentation()
    prs.slide_width = emu(1280)
    prs.slide_height = emu(720)
    blank = prs.slide_layouts[6]

    counts = {"rect": 0, "text": 0, "img": 0}
    for page in pages:
        slide = prs.slides.add_slide(blank)
        for s in page["shapes"]:
            k = s["k"]
            if k == "rect":
                add_rect(slide, s)
            elif k == "img":
                add_image(slide, s)
            elif k == "text":
                add_text(slide, s)
            counts[k] = counts.get(k, 0) + 1
        print(f"  p{page['page']:>2}: {len(page['shapes'])} 個の要素")

    prs.save(out_path)
    print(f"\n生成: {out_path}")
    print(f"  スライド {len(pages)} 枚 / 図形 {counts['rect']} ・ テキスト {counts['text']} ・ 画像 {counts['img']}")


if __name__ == "__main__":
    layout = sys.argv[1] if len(sys.argv) > 1 else "layout.json"
    out = sys.argv[2] if len(sys.argv) > 2 else "deck.pptx"
    build(layout, out)
